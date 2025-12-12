"""
PowerPoint (.pptx) Parser
Extracts text content and metadata from PowerPoint presentations
"""

import io
import logging
from typing import Dict, Any, List, Optional
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

class PowerPointParser:
    """Extract text and metadata from PowerPoint (.pptx) files"""
    
    def __init__(self):
        self.logger = logging.getLogger(__name__)
        
        if not PPTX_AVAILABLE:
            self.logger.warning("⚠️ python-pptx not available. Install with: pip install python-pptx")
    
    def can_process(self, file_path: str) -> bool:
        """Check if we can process this PowerPoint file"""
        if not PPTX_AVAILABLE:
            return False
        
        extension = Path(file_path).suffix.lower()
        return extension == '.pptx'
    
    def extract_text(self, content: bytes, file_path: str = "") -> str:
        """
        Extract all text content from PowerPoint presentation
        
        Args:
            content: PowerPoint file content as bytes
            file_path: File path for error reporting
            
        Returns:
            Extracted text content as string
        """
        if not PPTX_AVAILABLE:
            self.logger.error("❌ python-pptx not installed. Cannot process PowerPoint files.")
            return ""
        
        try:
            # Load presentation from bytes
            presentation_stream = io.BytesIO(content)
            presentation = Presentation(presentation_stream)
            
            extracted_text = []
            slide_count = 0
            
            # Extract text from all slides
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_count += 1
                slide_text = []
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Handle tables in slides
                    if hasattr(shape, "table"):
                        table_text = self._extract_table_text(shape.table)
                        if table_text:
                            slide_text.append(table_text)
                
                # Add slide content if any text found
                if slide_text:
                    slide_content = f"\\n--- Slide {slide_num} ---\\n" + "\\n".join(slide_text)
                    extracted_text.append(slide_content)
            
            # Combine all slide content
            full_text = "\\n\\n".join(extracted_text)
            
            # Add metadata summary
            if slide_count > 0:
                summary = f"PowerPoint Presentation Summary:\\nTotal Slides: {slide_count}\\nFile: {Path(file_path).name}\\n\\n"
                full_text = summary + full_text
            
            self.logger.debug(f"✅ Extracted text from {slide_count} slides in {Path(file_path).name}")
            
            return full_text if full_text.strip() else "PowerPoint presentation with no extractable text content."
            
        except PackageNotFoundError as e:
            self.logger.error(f"❌ Invalid PowerPoint file format: {file_path} - {e}")
            return f"Error: Invalid PowerPoint file format - {str(e)}"
            
        except Exception as e:
            self.logger.error(f"❌ Error extracting text from PowerPoint {file_path}: {e}")
            return f"Error extracting PowerPoint content: {str(e)}"
    
    def _extract_table_text(self, table) -> str:
        """Extract text from tables in PowerPoint slides"""
        try:
            table_text = []
            
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                
                if row_text:
                    table_text.append(" | ".join(row_text))
            
            return "\\n".join(table_text) if table_text else ""
            
        except Exception as e:
            self.logger.warning(f"⚠️ Error extracting table text: {e}")
            return ""
    
    def extract_metadata(self, content: bytes, file_path: str = "") -> Dict[str, Any]:
        """
        Extract metadata from PowerPoint presentation
        
        Args:
            content: PowerPoint file content as bytes
            file_path: File path for error reporting
            
        Returns:
            Dictionary with presentation metadata
        """
        if not PPTX_AVAILABLE:
            return {"error": "python-pptx not available"}
        
        try:
            presentation_stream = io.BytesIO(content)
            presentation = Presentation(presentation_stream)
            
            # Basic presentation metadata
            metadata = {
                "slide_count": len(presentation.slides),
                "file_type": "PowerPoint Presentation",
                "parser": "powerpoint_parser"
            }
            
            # Extract core properties if available
            try:
                core_props = presentation.core_properties
                
                if hasattr(core_props, 'title') and core_props.title:
                    metadata["title"] = core_props.title
                
                if hasattr(core_props, 'author') and core_props.author:
                    metadata["author"] = core_props.author
                
                if hasattr(core_props, 'subject') and core_props.subject:
                    metadata["subject"] = core_props.subject
                
                if hasattr(core_props, 'created') and core_props.created:
                    metadata["created"] = core_props.created.isoformat()
                
                if hasattr(core_props, 'modified') and core_props.modified:
                    metadata["modified"] = core_props.modified.isoformat()
                    
            except Exception as e:
                self.logger.warning(f"⚠️ Could not extract core properties: {e}")
            
            # Estimate content complexity
            total_shapes = 0
            slides_with_text = 0
            
            for slide in presentation.slides:
                shape_count = len(slide.shapes)
                total_shapes += shape_count
                
                # Check if slide has text content
                has_text = any(
                    hasattr(shape, "text") and shape.text.strip() 
                    for shape in slide.shapes
                )
                if has_text:
                    slides_with_text += 1
            
            metadata.update({
                "total_shapes": total_shapes,
                "slides_with_text": slides_with_text,
                "avg_shapes_per_slide": round(total_shapes / len(presentation.slides), 1) if presentation.slides else 0
            })
            
            return metadata
            
        except Exception as e:
            self.logger.error(f"❌ Error extracting PowerPoint metadata from {file_path}: {e}")
            return {
                "error": str(e),
                "file_type": "PowerPoint Presentation",
                "parser": "powerpoint_parser"
            }


# Test function
def test_powerpoint_parser():
    """Test the PowerPoint parser with a sample file"""
    parser = PowerPointParser()
    
    if not PPTX_AVAILABLE:
        print("❌ python-pptx not available. Install with: pip install python-pptx")
        return
    
    print("✅ PowerPoint parser ready")
    print(f"Can process .pptx files: {parser.can_process('test.pptx')}")


if __name__ == "__main__":
    test_powerpoint_parser()
