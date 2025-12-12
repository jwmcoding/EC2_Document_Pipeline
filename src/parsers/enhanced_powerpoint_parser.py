"""
Enhanced PowerPoint (.pptx) Parser with Vision Language Model Support
Extracts text content, images, and uses VLM to describe visual content
"""

import io
import base64
import logging
from typing import Dict, Any, List, Optional, Tuple
from pathlib import Path
import tempfile
import os

try:
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openai
    from PIL import Image
    VISION_AVAILABLE = True
except ImportError:
    VISION_AVAILABLE = False

class EnhancedPowerPointParser:
    """Enhanced PowerPoint parser with Vision Language Model support for image analysis"""
    
    def __init__(self, openai_client=None, enable_vision_analysis: bool = True, vision_model: str = "gpt-4o"):
        self.logger = logging.getLogger(__name__)
        self.openai_client = openai_client
        self.enable_vision_analysis = enable_vision_analysis and VISION_AVAILABLE
        self.vision_model = vision_model  # gpt-4o, gpt-4o-mini, or gpt-4-turbo
        
        if not PPTX_AVAILABLE:
            self.logger.warning("⚠️ python-pptx not available. Install with: pip install python-pptx")
        
        if not VISION_AVAILABLE and enable_vision_analysis:
            self.logger.warning("⚠️ Vision analysis dependencies missing. Install: pip install openai pillow")
    
    def can_process(self, file_path: str) -> bool:
        """Check if we can process this PowerPoint file"""
        if not PPTX_AVAILABLE:
            return False
        
        extension = Path(file_path).suffix.lower()
        return extension == '.pptx'
    
    def extract_enhanced_content(self, content: bytes, file_path: str = "") -> Dict[str, Any]:
        """
        Extract comprehensive content from PowerPoint including VLM image analysis
        
        Returns:
            Dictionary with text content, image descriptions, and metadata
        """
        if not PPTX_AVAILABLE:
            return {
                "text_content": "",
                "image_descriptions": [],
                "slide_count": 0,
                "error": "python-pptx not installed"
            }
        
        try:
            presentation_stream = io.BytesIO(content)
            presentation = Presentation(presentation_stream)
            
            result = {
                "text_content": "",
                "image_descriptions": [],
                "chart_descriptions": [],
                "diagram_descriptions": [],
                "slide_summaries": [],
                "slide_count": len(presentation.slides),
                "total_images": 0,
                "total_charts": 0,
                "processing_notes": []
            }
            
            all_text_parts = []
            
            # Process each slide
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_data = self._process_slide(slide, slide_num, file_path)
                
                # Accumulate slide content
                if slide_data["text_content"]:
                    all_text_parts.append(f"\\n--- Slide {slide_num} ---\\n{slide_data['text_content']}")
                
                # Accumulate image descriptions
                result["image_descriptions"].extend(slide_data["image_descriptions"])
                result["chart_descriptions"].extend(slide_data["chart_descriptions"])
                result["diagram_descriptions"].extend(slide_data["diagram_descriptions"])
                
                # Add slide summary
                result["slide_summaries"].append({
                    "slide_number": slide_num,
                    "text_length": len(slide_data["text_content"]),
                    "image_count": len(slide_data["image_descriptions"]),
                    "chart_count": len(slide_data["chart_descriptions"]),
                    "has_content": bool(slide_data["text_content"] or slide_data["image_descriptions"])
                })
                
                result["total_images"] += len(slide_data["image_descriptions"])
                result["total_charts"] += len(slide_data["chart_descriptions"])
            
            # Combine all text content
            if all_text_parts:
                summary = self._generate_presentation_summary(result)
                result["text_content"] = summary + "\\n\\n".join(all_text_parts)
            else:
                result["text_content"] = f"PowerPoint presentation: {Path(file_path).name} (no extractable content)"
            
            # Add enhanced descriptions to main text if available
            if result["image_descriptions"] or result["chart_descriptions"]:
                enhanced_content = self._combine_text_and_visual_analysis(result)
                result["text_content"] = enhanced_content
            
            self.logger.info(f"✅ Enhanced PowerPoint analysis: {result['slide_count']} slides, "
                           f"{result['total_images']} images, {result['total_charts']} charts")
            
            return result
            
        except Exception as e:
            self.logger.error(f"❌ Error in enhanced PowerPoint analysis {file_path}: {e}")
            return {
                "text_content": f"Error analyzing PowerPoint content: {str(e)}",
                "image_descriptions": [],
                "slide_count": 0,
                "error": str(e)
            }
    
    def _process_slide(self, slide, slide_num: int, file_path: str) -> Dict[str, Any]:
        """Process a single slide for text and visual content"""
        slide_data = {
            "text_content": "",
            "image_descriptions": [],
            "chart_descriptions": [],
            "diagram_descriptions": []
        }
        
        text_parts = []
        
        # Extract text and identify visual elements
        for shape in slide.shapes:
            # Extract text content
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
            
            # Handle tables (safe detection to avoid attribute errors)
            try:
                if getattr(shape, "has_table", False):
                    table_text = self._extract_table_text(shape.table)
                    if table_text:
                        text_parts.append(f"[Table]\n{table_text}")
            except Exception as e:
                self.logger.debug(f"Table extraction skipped on slide {slide_num}: {e}")
            
            # Process images and visual elements
            if self.enable_vision_analysis and self.openai_client:
                visual_description = self._analyze_shape_visual_content(shape, slide_num, file_path)
                if visual_description:
                    if "chart" in visual_description.lower() or "graph" in visual_description.lower():
                        slide_data["chart_descriptions"].append(visual_description)
                    elif "diagram" in visual_description.lower() or "flowchart" in visual_description.lower():
                        slide_data["diagram_descriptions"].append(visual_description)
                    else:
                        slide_data["image_descriptions"].append(visual_description)
        
        slide_data["text_content"] = "\\n".join(text_parts) if text_parts else ""
        return slide_data
    
    def _analyze_shape_visual_content(self, shape, slide_num: int, file_path: str) -> Optional[str]:
        """Analyze visual content of a shape using Vision Language Model"""
        try:
            # Check if shape contains visual content
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return self._analyze_image_with_vlm(shape, slide_num, "image")
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                return self._analyze_image_with_vlm(shape, slide_num, "chart")
            elif shape.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM]:
                # Could be a diagram or flowchart
                return self._analyze_image_with_vlm(shape, slide_num, "diagram")
            
            return None
            
        except Exception as e:
            self.logger.warning(f"⚠️ Could not analyze visual content on slide {slide_num}: {e}")
            return None
    
    def _analyze_image_with_vlm(self, shape, slide_num: int, content_type: str) -> Optional[str]:
        """Use Vision Language Model to analyze and describe visual content"""
        try:
            if not self.openai_client:
                return f"[{content_type.title()} on slide {slide_num} - Vision analysis not available]"

            # Attempt to extract an image for analysis (works for PICTURE shapes)
            image_base64 = self._extract_image_as_base64(shape, fmt="JPEG", max_dim=1024)

            if not image_base64:
                # If we cannot extract an image (e.g., CHART/DIAGRAM), provide a graceful fallback
                return f"[{content_type.title()} on slide {slide_num} - Image data not accessible for analysis]"

            # Call the Vision model with JPEG; on failure, try PNG fallback
            result_text = self._call_vision_model(image_base64=image_base64, content_type=content_type, slide_num=slide_num, mime="image/jpeg")
            if isinstance(result_text, str) and result_text.startswith("Vision analysis failed"):
                png_b64 = self._extract_image_as_base64(shape, fmt="PNG", max_dim=1024)
                if png_b64:
                    return self._call_vision_model(image_base64=png_b64, content_type=content_type, slide_num=slide_num, mime="image/png")
            return result_text
            
        except Exception as e:
            self.logger.warning(f"⚠️ VLM analysis failed for slide {slide_num}: {e}")
            return f"[{content_type.title()} on slide {slide_num} - Analysis failed: {str(e)}]"
    
    def _extract_image_as_base64(self, shape, fmt: str = "JPEG", max_dim: int = 1600) -> Optional[str]:
        """Extract image from PowerPoint shape and convert to base64 (re-encoded format).

        Notes:
        - Works for PICTURE shapes where python-pptx exposes shape.image.blob
        - For CHART/DIAGRAM shapes, python-pptx does not provide rasterized pixels; returns None
        - If Pillow is available, compress and resize large images for efficient API usage
        """
        try:
            # Picture shapes provide direct access to the embedded image
            if hasattr(shape, "image") and getattr(shape, "image") is not None:
                try:
                    image_bytes = shape.image.blob  # type: ignore[attr-defined]
                except Exception:
                    image_bytes = None

                if not image_bytes:
                    return None

                if not VISION_AVAILABLE:
                    return None
                try:
                    img = Image.open(io.BytesIO(image_bytes))
                    # Ensure RGB for consistent encoding
                    if img.mode not in ("RGB", "L"):
                        img = img.convert("RGB")
                    elif img.mode == "L":
                        img = img.convert("RGB")

                    # Resize if very large
                    width, height = img.size
                    longest = max(width, height)
                    if longest > max_dim:
                        scale = max_dim / float(longest)
                        new_size = (int(width * scale), int(height * scale))
                        img = img.resize(new_size)

                    buffer = io.BytesIO()
                    save_kwargs = {"optimize": True}
                    if fmt.upper() == "JPEG":
                        save_kwargs.update({"quality": 85})
                    img.save(buffer, format=fmt.upper(), **save_kwargs)
                    return base64.b64encode(buffer.getvalue()).decode("utf-8")
                except Exception as e:
                    # If re-encode fails, do not send original bytes
                    self.logger.debug(f"Pillow re-encode failed; skipping image analysis: {e}")
                    return None

            # No accessible image data for this shape type
            return None
        except Exception as e:
            self.logger.warning(f"⚠️ Could not extract image for VLM analysis: {e}")
            return None
    
    def _call_vision_model(self, image_base64: str, content_type: str, slide_num: int, mime: str = "image/jpeg") -> str:
        """Call OpenAI Vision API to analyze image content following latest OpenAI guidelines"""
        try:
            prompt = self._create_vision_prompt(content_type, slide_num)
            
            response = self.openai_client.chat.completions.create(
                model=self.vision_model,  # Configurable vision model
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:{mime};base64,{image_base64}",
                                    "detail": "high"
                                }
                            }
                        ]
                    }
                ],
                max_tokens=500,
                temperature=0.1
            )
            
            return response.choices[0].message.content
        except Exception as e:
            self.logger.error(f"❌ Vision API call failed: {e}")
            self.logger.debug(f"Image size: {len(image_base64)} chars, Content type: {content_type}")
            return f"Vision analysis failed: {str(e)}"
    
    def _create_vision_prompt(self, content_type: str, slide_num: int) -> str:
        """Create specialized prompt for different types of visual content"""
        
        if content_type == "chart":
            return f"""Analyze this chart/graph from slide {slide_num} of a business presentation. Provide:
1. Chart type (bar, line, pie, etc.)
2. Key data insights and trends
3. Main message or conclusion
4. Business relevance
Keep response concise and business-focused."""
        
        elif content_type == "diagram":
            return f"""Analyze this diagram/flowchart from slide {slide_num} of a business presentation. Describe:
1. Type of diagram (process flow, org chart, system architecture, etc.)
2. Key components and relationships
3. Main process or concept illustrated
4. Business purpose or insight
Keep response concise and business-focused."""
        
        else:  # general image
            return f"""Analyze this image from slide {slide_num} of a business presentation. Describe:
1. What the image shows
2. Key visual elements
3. Business context or relevance
4. Any text or data visible
Keep response concise and business-focused."""
    
    def _extract_table_text(self, table) -> str:
        """Extract text from tables in PowerPoint slides"""
        try:
            table_text = []
            
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                
                if row_text:
                    table_text.append(" | ".join(row_text))
            
            return "\\n".join(table_text) if table_text else ""
            
        except Exception as e:
            self.logger.warning(f"⚠️ Error extracting table text: {e}")
            return ""
    
    def _generate_presentation_summary(self, result: Dict[str, Any]) -> str:
        """Generate a summary of the presentation content"""
        summary_parts = [
            f"PowerPoint Presentation Analysis Summary:",
            f"Total Slides: {result['slide_count']}",
            f"Images Analyzed: {result['total_images']}",
            f"Charts/Graphs: {result['total_charts']}"
        ]
        
        if result["image_descriptions"]:
            summary_parts.append(f"Visual Content: Enhanced with AI-generated descriptions")
        
        return "\\n".join(summary_parts) + "\\n\\n"
    
    def _combine_text_and_visual_analysis(self, result: Dict[str, Any]) -> str:
        """Combine text content with visual analysis for comprehensive understanding"""
        content_parts = [result["text_content"]]
        
        if result["image_descriptions"]:
            content_parts.append("\\n--- Visual Content Analysis ---")
            content_parts.extend([f"Image: {desc}" for desc in result["image_descriptions"]])
        
        if result["chart_descriptions"]:
            content_parts.append("\\n--- Chart/Graph Analysis ---")
            content_parts.extend([f"Chart: {desc}" for desc in result["chart_descriptions"]])
        
        if result["diagram_descriptions"]:
            content_parts.append("\\n--- Diagram Analysis ---")
            content_parts.extend([f"Diagram: {desc}" for desc in result["diagram_descriptions"]])
        
        return "\\n".join(content_parts)
    
    def extract_text(self, content: bytes, file_path: str = "") -> str:
        """Extract text content (compatibility method)"""
        result = self.extract_enhanced_content(content, file_path)
        return result.get("text_content", "")


# Factory function for easy integration
def create_powerpoint_parser(openai_client=None, enable_vision: bool = True, vision_model: str = "gpt-4o"):
    """Create enhanced PowerPoint parser with optional VLM capabilities
    
    Args:
        openai_client: OpenAI client instance
        enable_vision: Enable VLM analysis
        vision_model: Vision model to use ("gpt-4o", "gpt-4o-mini", "gpt-4-turbo")
    """
    return EnhancedPowerPointParser(
        openai_client=openai_client, 
        enable_vision_analysis=enable_vision,
        vision_model=vision_model
    )
